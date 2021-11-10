import random
import os
from wand.image import Image
from pptx import Presentation
directory_path = os.path.dirname(os.path.realpath(__file__))               # script directory path

# ----------------------FUNCTION TO ADD WATERMARK TO IMAGES-------------------------------------------


def add_watermark_to_image():
    # looping to add watermark to each image for a given range
    for i in range(1, 6):
        # importing the image to be watermarked
        with Image(filename=f"{directory_path}\image{i}.jpg") as image:

            # importing the nike logo image
            with Image(filename=f'{directory_path}\\nike_black.png') as nike_logo:
                # Clone the image needs to be watermarked in order to process
                with image.clone() as image_to_be_watermark:
                    # left as 10 and top as 10 for proper adjustment of logo
                    nike_logo.resize(1300, 300)
                    image_to_be_watermark.watermark(nike_logo, 0, 10, 10)

                    # Saving the output image
                    image_to_be_watermark.save(
                        filename=f"{directory_path}\watermark_image{i}.jpg")

# -----------FUNCTION TO INSERT IMAGES IN PPT AND SAVE THE PPT--------------------------------

def add_image_to_ppt(image_url):
    from PIL import Image
    layout = prs.slide_layouts[8]
    slide = prs.slides.add_slide(layout)
    # adding title to slide
    title = slide.shapes.title.text = f"Slide {i}"
    # adding subtitle to slide
    subtitle = slide.placeholders[2].text = f"Image {i}"
    placeholder = slide.placeholders[1]

    # height and width dimensions of the opened image which needs to be added further
    image_file = Image.open(image_url)
    width, height = image_file.size

    placeholder.height = height
    placeholder.width = width

    # Inserting image in placeholder of slide
    placeholder = placeholder.insert_picture(image_url)

    # Calculating ratios for comparisons
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio

    #  conditon for Placeholder width if it is too wide:
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    #  or if Placeholder height is too high then
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side


# ------------------D R I V E R C O D E-----------------------------
prs = Presentation()
# callling function for adding watermark to images
add_watermark_to_image()
# looping to add images to ppt slides in a given range
for i in range(1, 6):
    # calling function to add images in ppt slides
    add_image_to_ppt(f"{directory_path}\watermark_image{i}.jpg")
    # removing watermarked images from directory after ppt creation
    os.remove(f"{directory_path}\watermark_image{i}.jpg")

# using random module to save ppt with unique name
random_number = random.randint(100000, 9999999)
# saving ppt to the directory
prs.save(f"{directory_path}\\file_{random_number}.pptx")





